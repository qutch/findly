"""File parsers — PDF, text, code extraction."""
import pymupdf
import os
from docx import Document
from pptx import Presentation
from datetime import datetime

class File:
    def __init__(self, fileName: str):
        parsed = FileProcessor.parseFile(fileName)
        self.metadata = parsed["metadata"] # Dict
        self.content = parsed["content"] # String

class FileProcessor:
    # Empty constructor since all methods are static, but allows for future state if needed
    def __init__(self):
        pass

    # PDF PARSER
    @staticmethod
    def parsePdf(fileName: str) -> str:
        doc = pymupdf.open(fileName)
        try:
            text = ""
            for page in doc:
                text += page.get_text()
            return text
        finally:
            doc.close()

    # TEXT + CODE PARSER
    @staticmethod
    def parseTxt(fileName: str) -> str:
        with open(fileName, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    # DOCX PARSER
    @staticmethod
    def parseDocx(fileName: str) -> str:
        doc = Document(fileName)
        return "\n".join(p.text for p in doc.paragraphs)

    # PPTX PARSER
    @staticmethod
    def parsePptx(fileName: str) -> str:
        pres = Presentation(fileName)
        parts = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    parts.append(shape.text)
        return "\n".join(parts)

    # Format bytes to human-readable format
    @staticmethod
    def format_bytes(bytes_size: int) -> str:
        """Convert bytes to human-readable format."""
        for unit in ['B', 'KB', 'MB', 'GB', 'TB']:
            if bytes_size < 1024.0:
                return f"{bytes_size:.1f} {unit}"
            bytes_size /= 1024.0
        return f"{bytes_size:.1f} PB"
    

    # Extract metadata such as filename, type, size, and timestamps
    @staticmethod
    def extractMetadata(fileName: str) -> dict:
        """
        Extract file metadata in a format compatible with ranking.py.
        Returns camelCase keys to match ranking service expectations.
        """
        stats = os.stat(fileName)

        return {
            # camelCase for compatibility with ranking.py
            "fileName": os.path.basename(fileName),
            "filePath": os.path.abspath(fileName),
            "fileType": os.path.splitext(fileName)[1].lower(),
            "fileSize": stats.st_size,  # Size in bytes
            "sizeReadable": FileProcessor.format_bytes(stats.st_size), # Human-readable size (2.5 MB)

            # Timestamps as raw values
            "lastModified": stats.st_mtime,
            "lastAccessed": stats.st_atime,

            # Timestamps as readable strings (ISO format for ranking.py)
            "lastModifiedReadable": datetime.fromtimestamp(stats.st_mtime).isoformat(),
            "lastAccessedReadable": datetime.fromtimestamp(stats.st_atime).isoformat(),
        }

    @staticmethod
    def parseFile(fileName: str) -> dict:
        """
        Parse a file and extract both metadata and content.

        Handles case-insensitive file extensions and files without extensions.
        """
        # Extracts metadata
        metadata = FileProcessor.extractMetadata(fileName)

        # Parses the file content - normalize to lowercase for case-insensitive matching
        file_type = metadata["fileType"].lower()
        content = ""

        if file_type == ".pdf":
            content = FileProcessor.parsePdf(fileName)
        elif file_type in (".txt", ".md", ".py", ".js", ".ts", ".json", ".csv", ".html", ".css", ""):
            # Handle plain-text files: text, markdown, code, data, and extensionless (README, Makefile, etc.)
            content = FileProcessor.parseTxt(fileName)
        elif file_type == ".docx":
            content = FileProcessor.parseDocx(fileName)
        elif file_type == ".pptx":
            content = FileProcessor.parsePptx(fileName)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

        return {
            "metadata": metadata,
            "content": content
        }


    # Takes in list of filenames from pinecone service
    @staticmethod
    def sendToRankingService(fileNames: list[str]) -> list[File]:
        """
        Send files to ranking service for processing and return ranked list of File objects.
        """
        files = []
        for fileName in fileNames:
            files.append(File(fileName))

        return files


    @staticmethod
    def chunkContent(content: str, chunk_size: int = 500, overlap: int = 100) -> list[str]:
        """
        Chunk content into pieces with optional overlap for better context in vectorization.
        """
        if not content or len(content) <= chunk_size:
            return [content] if content else []

        chunks = []
        start = 0

        while start < len(content):
            # Get chunk from start to start + chunk_size
            end = start + chunk_size
            chunk = content[start:end]

            # If not the last chunk and we're not at a natural break, try to break at sentence/word
            if end < len(content):
                # Try to break at sentence (period, question mark, exclamation)
                last_sentence = max(chunk.rfind('. '), chunk.rfind('? '), chunk.rfind('! '))
                if last_sentence > chunk_size * 0.5:  # Only break at sentence if it's past halfway
                    chunk = chunk[:last_sentence + 1]
                    end = start + last_sentence + 1
                else:
                    # Otherwise try to break at word boundary
                    last_space = chunk.rfind(' ')
                    if last_space > 0:
                        chunk = chunk[:last_space]
                        end = start + last_space

            chunks.append(chunk.strip())

            # Move start forward, accounting for overlap
            start = end - overlap if end < len(content) else end

        return chunks


    @staticmethod
    def prepareForPinecone(fileName: str, chunk_size: int = 500, overlap: int = 100) -> dict:
        """
        Parse a file and prepare it for Pinecone ingestion.
        """
        # Parse file to get content and metadata
        parsed = FileProcessor.parseFile(fileName)

        # Chunk the content
        chunks = FileProcessor.chunkContent(parsed["content"], chunk_size, overlap)

        return {
            "chunks": chunks,
            "metadata": parsed["metadata"],
        }
    
    @staticmethod
    def sendToPinecone(fileName: str):
        """
        Full pipeline to parse a file, prepare it, and send it to Pinecone.
        """
        preparedFiles = FileProcessor.prepareForPinecone(fileName)
        print(fileName)
        print(preparedFiles)
        
        # Call Pinecone Service to upsert chunks and metadata (not implemented here)